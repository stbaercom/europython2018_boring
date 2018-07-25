import pandas as pd
import pptx

from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pd2ppt import df_to_table


def load_excel_files():
    df_times = pd.read_excel("input_data/project_hours.xlsx")
    df_expenses = pd.read_excel("input_data/project_expenses.xlsx")
    df_rates = pd.read_excel("input_data/project_rates.xlsx")
    return df_times, df_expenses, df_rates


def transform_excel(df_times, df_expenses, df_rates):
    df_times_rate = df_times.merge(df_rates, how="outer", on="Person")
    times_diff = df_times_rate["TimeStop"] - df_times_rate["TimeStart"]
    df_times_rate["Cost"] = times_diff * df_times_rate["Rate"]
    df_times_cost_pivot = df_times_rate.pivot_table(
        values="Cost", index=["Project", "Person"]).reset_index()
    df_times_cost_pivot["Cost Type"] = "hours"
    df_expenses_pivot = df_expenses.pivot_table(
        values="Cost", index=["Project", "Person"]).reset_index()
    df_expenses_pivot["Cost Type"] = "expenses"
    df_all_costs = pd.concat([df_expenses_pivot, df_times_cost_pivot], sort=True)
    return df_times_cost_pivot, df_expenses_pivot, df_all_costs


def create_introsheet(workbook):
    introsheet = workbook.add_worksheet("Introduction")
    bold = workbook.add_format(
        {'bold': True, "align": "right", "font_color": "blue"})
    introsheet.write(0, 0, 'Title', bold)
    intro_text = 'Overall Costs'
    introsheet.write(0, 1, 'Overall Costs')
    introsheet.set_column(1, 1, len(intro_text) + 5)
    introsheet.insert_image(1, 0, "input_data/logo.jpg",
                            {'x_scale': 0.5, 'y_scale': 0.5})


def export_to_xlsx_sheets(df_times_cost_pivot, df_expenses_pivot, df_all_costs):
    writer = pd.ExcelWriter('scrap_data/pandas_simple.xlsx')
    df_all_costs.to_excel(writer, index=False, sheet_name='df_all_costs')
    df_expenses_pivot.to_excel(writer, index=False, sheet_name='df_expenses_pivot')
    df_times_cost_pivot.to_excel(writer, index=False, sheet_name='df_times_cost_pivot')
    writer.close()


def create_sheets_from_pandas_intro(df_times_cost_pivot, df_expenses_pivot, df_all_costs):
    writer = pd.ExcelWriter('scrap_data/pandas_simple_intro.xlsx',
                            engine='xlsxwriter')
    workbook = writer.book
    create_introsheet(workbook)
    df_all_costs.to_excel(writer, index=False,
                          sheet_name='df_all_costs')
    df_expenses_pivot.to_excel(writer, index=False,
                               sheet_name='df_expenses_pivot')
    df_times_cost_pivot.to_excel(writer, index=False,
                                 sheet_name='df_times_cost_pivot')


def create_pandas_by_hand_1(workbook, sheet_title, dataframe):
    sheet = workbook.add_worksheet(sheet_title)
    sheet.write_row(0, 0, dataframe.columns)
    for i, row in enumerate(dataframe.values):
        sheet.write_row(i + 1, 0, row)


def create_pandas_by_hand_2(workbook, sheet_title, dataframe):
    sheet = workbook.add_worksheet(sheet_title)
    large_text = workbook.add_format({'bold': True, "font_size": 14})
    red_bold = workbook.add_format({'bold': True, "font_color": "red"})
    sheet.write_row(0, 0, dataframe.columns, large_text)

    for i, header in enumerate(dataframe.columns):
        sheet.set_column(i, i, len(header) * 1.2 + 5)

    percentile75 = dataframe["Cost"].describe()["75%"]
    for i, row in enumerate(dataframe.values):
        for i2, value in enumerate(row):
            if i2 == 0:
                if value > percentile75:
                    sheet.write_number(i + 1, i2, value, red_bold)
                else:
                    sheet.write_number(i + 1, i2, value)
            else:
                sheet.write_string(i + 1, i2, value)


def create_pandas_by_hand_3(workbook, sheet_title, dataframe):
    num_format = workbook.add_format({'num_format': "####.#"})
    sheet = workbook.add_worksheet(sheet_title)
    nrows, ncols = dataframe.shape
    columns_desc = [{"header": v} for v in dataframe.columns]
    sheet.add_table(0, 0, nrows, ncols - 1, {"data": dataframe.values,
                                             "columns": columns_desc})
    sheet.set_column(0, 0, 10, num_format)

    conditional_options = {
        'type': '3_color_scale',
        "min_color": "green",
        "mid_color": "yellow",
        "max_color": "red"
    }
    sheet.conditional_format(1, 0, nrows, 0, conditional_options)


def create_chart_1(workbook, sheet_title, df_all_costs):
    sheet = workbook.add_worksheet(sheet_title)
    df_chart = df_all_costs.pivot_table(
        values="Cost", index="Person", columns="Cost Type")
    df_chart.reset_index(inplace=True)
    sheet.write_row(0, 0, [s.upper() for s in df_chart.columns])
    sheet.write_column(1, 0, df_chart['Person'])
    sheet.write_column(1, 1, df_chart['expenses'])
    sheet.write_column(1, 2, df_chart['hours'])
    chart = workbook.add_chart({'type': 'column', 'subtype': 'stacked'})
    chart.set_style(12)
    nrows = df_chart.shape[0]
    for i in [1, 2]:
        chart.add_series({
            'name': [sheet.get_name(), 0, i],
            'categories': [sheet.get_name(), 1, 0, nrows, 0],
            'values': [sheet.get_name(), 1, i, nrows, i]})
    sheet.insert_chart('A8', chart, {'x_offset': 25, 'y_offset': 10})


def prepare_excel_xlsxwriter(df_all_costs, df_expenses_pivot, df_times_cost_pivot):
    export_to_xlsx_sheets(df_times_cost_pivot, df_expenses_pivot, df_all_costs)
    create_sheets_from_pandas_intro(df_times_cost_pivot, df_expenses_pivot, df_all_costs)
    writer = pd.ExcelWriter('scrap_data/pandas_complex.xlsx', engine='xlsxwriter')
    workbook = writer.book
    create_pandas_by_hand_1(workbook, "All Costs", df_all_costs)
    create_pandas_by_hand_2(workbook, "All Costs 2", df_all_costs)
    create_pandas_by_hand_3(workbook, "All Costs 3", df_all_costs)
    create_chart_1(workbook, "Sheet with Chart 1", df_all_costs)


def create_slide(presentation, title, layout=5):
    layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(layout)
    if title is not None:
        slide_title = slide.shapes.title
        slide_title.text = title
    return slide


def prepare_pptx(df_all_costs):
    create_presentation_1()

    presentation = pptx.Presentation("input_data/template.pptx")

    slide = create_slide(presentation, "Introduction")
    create_intro_slide_with_graphic(slide)

    slide = create_slide(presentation, "Data Table")
    create_table_slide(df_all_costs, slide)

    slide = create_slide(presentation, "Charts")

    create_chart_slide(df_all_costs, slide)

    presentation.save('./output_data/test.pptx')


def prepare_pptx_and_convert(df_all_costs, pptx_filename, export_format="pdf"):
    import os
    import subprocess
    import pptx

    presentation_plain = pptx.Presentation("input_data/template_plain.pptx")
    slide = create_slide(presentation_plain, "Charts")
    create_chart_slide(df_all_costs, slide)
    presentation_plain.save(pptx_filename)

    libre_office_binary = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    cmd = [libre_office_binary, "--headless", "--convert-to", export_format,
           "--outdir", os.path.dirname(pptx_filename),
           pptx_filename]
    subprocess.run(cmd, check=True)


def combine_pdf(pptx_filename):
    import pdfrw
    pdf_filename = pptx_filename.replace(".pptx", ".pdf")
    pdf_report_pages = pdfrw.PdfReader(pdf_filename).pages
    pdf_template_pages = pdfrw.PdfReader('input_data/pdf_template.pdf').pages
    outdata = pdfrw.PdfWriter('output_data/plain_with_template.pdf')
    outdata.addpage(pdf_template_pages[0])
    outdata.addpages(pdf_report_pages)
    outdata.addpage(pdf_template_pages[1])
    outdata.write()


def create_chart_slide(df_all_costs, slide):
    df_chart = df_all_costs.pivot_table(values="Cost",
                                        index="Person", columns="Cost Type")
    df_chart.reset_index(inplace=True)
    chart_data = ChartData()
    chart_data.categories = list(df_chart['Person'])
    chart_data.add_series('Expenses', list(df_chart["expenses"]))
    chart_data.add_series('Hours', list(df_chart["hours"]))
    CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED
    chart_left = Inches(1);
    chart_top = Inches(2)
    chart_width = Inches(12);
    chart_height = Inches(4)
    chart = slide.shapes.add_chart(CHART_TYPE, chart_left, chart_top,
                                   chart_width, chart_height, chart_data).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False


def create_table_slide(df_all_costs, slide):
    table_left = Inches(1);
    table_top = Inches(2)
    table_width = Inches(12);
    table_height = Inches(4)
    df_to_table(slide, df_all_costs, table_left, table_top,
                table_width, table_height)


def create_intro_slide_with_graphic(slide):
    left = width = height = Inches(1)
    top = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "A Short but meaningful text for the slide"
    top = Inches(4)
    slide.shapes.add_picture("./input_data/logo.jpg", left, top)


def create_presentation_1():
    presentation = pptx.Presentation("input_data/template.pptx")
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Meaningful Title"
    subtitle = slide.placeholders[1]
    subtitle.text = "Some text for the placeholder defined in the layout"
    presentation.save("./output_data/presentation_1.pptx")


def prepare_pdf(df_all_costs):
    pptx_filename = './output_data/plain.pptx'
    prepare_pptx_and_convert(df_all_costs, pptx_filename)
    combine_pdf(pptx_filename)


def main():
    # Just load the data from Excel files and rename some columns
    df_times, df_expenses, df_rates = load_excel_files()

    # Build some Pivot tables, because everybody _loves_ pivot tables
    df_times_cost_pivot, df_expenses_pivot, df_all_costs = transform_excel(df_times, df_expenses, df_rates)

    # Create the different versions of Excel file, in increasing order of colorfulness...
    prepare_excel_xlsxwriter(df_all_costs, df_expenses_pivot, df_times_cost_pivot)

    # Prepare a PPTX, based on the pivots and an existing PPTX 'template'
    prepare_pptx(df_all_costs)

    # Finally, create a version of the PPTX to turn into a PDF via Libreoffice, and process the resulting file
    # with Python
    prepare_pdf(df_all_costs)


if __name__ == '__main__':
    main()
